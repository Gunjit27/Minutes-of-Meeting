import json
import re
from typing import List
from langchain_ollama import ChatOllama
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import state


# ---------------- LLM SETUP ---------------- #

llm = ChatOllama(
    model="llama3.1:8b",
    temperature=0.1
)

prompt = """
You are an information extraction system.

STRICT RULES:
- ONLY extract information explicitly present in transcript
- DO NOT infer, assume, summarize creatively, or rephrase meaning
- DO NOT add new numbers or modify values
- If unsure → SKIP
- If section missing → return []

- Every bullet MUST be directly traceable to transcript text

OUTPUT MUST BE VALID JSON ONLY

JSON SCHEMA:
{{
  "title": "string",
  "slides": [
    {{
      "title": "string",
      "bullets": ["string"]
    }}
  ],
  "charts": []
}}

IMPORTANT:
- Only include charts if EXACT numeric comparisons exist
- DO NOT invent charts

TRANSCRIPT:
{transcript}
"""


# ---------------- JSON HELPERS ---------------- #

def extract_json(text: str):
    """Extract JSON block from messy LLM output"""
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        return match.group(0)
    raise ValueError("No JSON found")


def generate_json(transcript: str, retries=3):
    for attempt in range(retries):
        response = llm.invoke(prompt.format(transcript=transcript))
        raw = response.content.strip()

        try:
            json_str = extract_json(raw)
            return json.loads(json_str)

        except Exception as e:
            print(f"[Retry {attempt+1}] JSON parsing failed:", e)

    raise ValueError("Failed to generate valid JSON")


# ---------------- PPT GENERATION ---------------- #

def add_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for item in content_list:
        p = tf.add_paragraph()
        p.text = str(item)
        p.level = 0


def add_chart_slide(prs, chart):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = chart["title"]

    chart_data = ChartData()
    chart_data.categories = chart["labels"]
    chart_data.add_series("Series 1", chart["values"])

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(6), Inches(4),
        chart_data
    )


def create_ppt(data, filename="output.pptx"):
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = data.get("title", "Meeting Summary")

    # Slides (FIXED — uses correct schema)
    for slide_data in data.get("slides", []):
        add_slide(
            prs,
            slide_data.get("title", ""),
            slide_data.get("bullets", [])
        )

    # Charts
    for chart in data.get("charts", []):
        if chart.get("labels") and chart.get("values"):
            add_chart_slide(prs, chart)

    prs.save(filename)
    print(f"✅ PPT saved as {filename}")


# ---------------- MAIN ---------------- #

def run_ppt_generation():

    transcript = state.get_transcript()
    if not transcript:
        raise ValueError("Transcript not available. Please upload audio first.")

    data = generate_json(transcript)

    filename = "output.pptx"
    create_ppt(data, filename)

    return filename